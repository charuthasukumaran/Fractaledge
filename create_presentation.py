"""
Generate an attractive PowerPoint presentation for NIFTY Stress Monitor.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Color Palette (Dark Professional Theme) ──────────────────────────
BG_DARK      = RGBColor(0x0F, 0x17, 0x29)   # Deep navy background
BG_CARD      = RGBColor(0x1A, 0x25, 0x3C)   # Card background
ACCENT_BLUE  = RGBColor(0x38, 0xBD, 0xF8)   # Bright cyan accent
ACCENT_GREEN = RGBColor(0x22, 0xC5, 0x5E)   # Green
ACCENT_AMBER = RGBColor(0xF5, 0xA6, 0x23)   # Amber/Orange
ACCENT_RED   = RGBColor(0xEF, 0x44, 0x44)   # Red
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xCB, 0xD5, 0xE1)   # Subtle text
MID_GRAY     = RGBColor(0x94, 0xA3, 0xB8)
PURPLE       = RGBColor(0xA7, 0x8B, 0xFA)   # Purple accent
GRADIENT_TOP = RGBColor(0x0F, 0x17, 0x29)
GRADIENT_BOT = RGBColor(0x1E, 0x29, 0x3B)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


# ── Helper Functions ──────────────────────────────────────────────────

def add_bg(slide, color=BG_DARK):
    """Fill slide background with solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, border_color=None, radius=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color or BG_CARD
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE,
                 bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    """Add a text box with styled text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multi_text(slide, left, top, width, height, lines, default_size=16,
                   default_color=LIGHT_GRAY, line_spacing=1.3):
    """Add text box with multiple styled lines. Each line: (text, size, color, bold)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        text = line_data[0]
        size = line_data[1] if len(line_data) > 1 else default_size
        color = line_data[2] if len(line_data) > 2 else default_color
        bold = line_data[3] if len(line_data) > 3 else False
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "Segoe UI"
        p.space_after = Pt(size * (line_spacing - 1) + 2)
    return txBox


def add_accent_line(slide, left, top, width, color=ACCENT_BLUE):
    """Add a thin accent line."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_circle(slide, left, top, size, fill_color):
    """Add a circle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_icon_card(slide, left, top, width, height, icon_text, title, desc,
                  icon_color=ACCENT_BLUE, title_color=WHITE, desc_color=LIGHT_GRAY):
    """Add a card with an icon circle, title, and description."""
    card = add_shape(slide, left, top, width, height, BG_CARD, border_color=RGBColor(0x2D, 0x3A, 0x55))
    # Icon circle
    circle = add_circle(slide, left + Inches(0.3), top + Inches(0.35), Inches(0.6), icon_color)
    add_text_box(slide, left + Inches(0.3), top + Inches(0.35), Inches(0.6), Inches(0.6),
                 icon_text, font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Title
    add_text_box(slide, left + Inches(1.1), top + Inches(0.25), width - Inches(1.4), Inches(0.4),
                 title, font_size=16, color=title_color, bold=True)
    # Description
    add_text_box(slide, left + Inches(1.1), top + Inches(0.65), width - Inches(1.4), height - Inches(0.8),
                 desc, font_size=12, color=desc_color)


def slide_header(slide, number, title, subtitle=""):
    """Add a standard slide header with number badge and accent line."""
    add_bg(slide)
    # Slide number badge
    badge = add_shape(slide, Inches(0.6), Inches(0.4), Inches(0.55), Inches(0.55), ACCENT_BLUE)
    add_text_box(slide, Inches(0.6), Inches(0.42), Inches(0.55), Inches(0.55),
                 str(number), font_size=20, color=BG_DARK, bold=True, alignment=PP_ALIGN.CENTER)
    # Title
    add_text_box(slide, Inches(1.35), Inches(0.35), Inches(10), Inches(0.6),
                 title, font_size=32, color=WHITE, bold=True)
    # Accent line
    add_accent_line(slide, Inches(0.6), Inches(1.1), Inches(2), ACCENT_BLUE)
    if subtitle:
        add_text_box(slide, Inches(0.6), Inches(1.25), Inches(11), Inches(0.45),
                     subtitle, font_size=16, color=MID_GRAY)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 1: TITLE SLIDE
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)

# Large decorative circles (background elements)
add_circle(slide, Inches(-1), Inches(-1), Inches(4), RGBColor(0x1A, 0x25, 0x3C))
add_circle(slide, Inches(10), Inches(4.5), Inches(5), RGBColor(0x15, 0x20, 0x35))

# Top accent line
add_accent_line(slide, Inches(1.5), Inches(1.5), Inches(1.2), ACCENT_BLUE)

# Main title
add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.2),
             "NIFTY Stress Monitor", font_size=52, color=WHITE, bold=True)

# Subtitle
add_text_box(slide, Inches(1.5), Inches(3.0), Inches(10), Inches(0.6),
             "Real-Time Market Stress Detection & Trading Analysis Platform",
             font_size=22, color=ACCENT_BLUE)

# Description
add_text_box(slide, Inches(1.5), Inches(3.8), Inches(8), Inches(1),
             "Powered by Multifractal Mathematics, AI Insights & Real-Time Indian Market Data",
             font_size=16, color=MID_GRAY)

# Three small feature pills at bottom
pill_y = Inches(5.2)
pills = [("Fractal Analysis", ACCENT_BLUE), ("70+ Stocks", ACCENT_GREEN), ("AI-Powered", PURPLE)]
for i, (text, color) in enumerate(pills):
    x = Inches(1.5 + i * 2.8)
    pill = add_shape(slide, x, pill_y, Inches(2.4), Inches(0.5), color)
    add_text_box(slide, x, pill_y + Inches(0.05), Inches(2.4), Inches(0.4),
                 text, font_size=14, color=BG_DARK, bold=True, alignment=PP_ALIGN.CENTER)

# Bottom tagline
add_text_box(slide, Inches(1.5), Inches(6.3), Inches(10), Inches(0.4),
             "Built with Python  |  FastAPI  |  NumPy  |  React  |  Claude AI",
             font_size=13, color=MID_GRAY, alignment=PP_ALIGN.LEFT)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 2: THE PROBLEM
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, 1, "The Problem We Solve", "Why traditional tools fall short")

# Left column - Problem cards
problems = [
    ("Lagging Indicators", "RSI, MACD, and Moving Averages tell you what\nALREADY happened. Like driving with a rearview mirror."),
    ("Regime Blindness", "Traditional tools use the same logic in calm AND\nchaotic markets. They don't detect regime shifts."),
    ("False Signals", "Fixed-parameter indicators produce excessive false\nsignals when market character changes."),
    ("No Cross-Asset View", "Single-stock analysis misses critical sector\ncoupling and divergence signals."),
]
for i, (title, desc) in enumerate(problems):
    y = Inches(1.85 + i * 1.25)
    add_shape(slide, Inches(0.6), y, Inches(5.8), Inches(1.1), BG_CARD, RGBColor(0x2D, 0x3A, 0x55))
    # Red X circle
    add_circle(slide, Inches(0.85), y + Inches(0.25), Inches(0.5), ACCENT_RED)
    add_text_box(slide, Inches(0.85), y + Inches(0.25), Inches(0.5), Inches(0.5),
                 "X", font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.55), y + Inches(0.12), Inches(4.6), Inches(0.35),
                 title, font_size=16, color=WHITE, bold=True)
    add_text_box(slide, Inches(1.55), y + Inches(0.48), Inches(4.6), Inches(0.6),
                 desc, font_size=12, color=LIGHT_GRAY)

# Right column - Solution
add_shape(slide, Inches(6.8), Inches(1.85), Inches(5.9), Inches(4.9), RGBColor(0x0D, 0x2A, 0x1A), RGBColor(0x16, 0x65, 0x34))
add_text_box(slide, Inches(7.1), Inches(2.05), Inches(5.3), Inches(0.5),
             "OUR SOLUTION", font_size=14, color=ACCENT_GREEN, bold=True)
add_accent_line(slide, Inches(7.1), Inches(2.55), Inches(1.5), ACCENT_GREEN)
solutions = [
    "Detects market regime shifts BEFORE they become obvious using multifractal mathematics",
    "Adapts all analysis (stops, targets, quality scores) to the current market regime",
    "Monitors cross-asset coupling (NIFTY vs BANKNIFTY) for early stress warnings",
    "Combines fractal science with AI to deliver plain-English actionable insights",
    "Works with free data - no paid subscriptions required",
]
for i, sol in enumerate(solutions):
    y = Inches(2.8 + i * 0.75)
    add_circle(slide, Inches(7.2), y + Inches(0.05), Inches(0.3), ACCENT_GREEN)
    add_text_box(slide, Inches(7.2), y + Inches(0.02), Inches(0.3), Inches(0.3),
                 "\u2713", font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(7.7), y, Inches(4.7), Inches(0.65),
                 sol, font_size=13, color=LIGHT_GRAY)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 3: HOW IT WORKS (SIMPLE)
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, 2, "How It Works", "A simple explanation anyone can understand")

# The analogy
add_text_box(slide, Inches(0.6), Inches(1.6), Inches(12), Inches(0.5),
             "Think of the stock market like an ocean...", font_size=20, color=ACCENT_BLUE, bold=True)

steps = [
    ("\U0001F30A", "Watch the Waves", "Collects price data every 5 minutes.\nLooks at ~1000 price points\n(3-4 trading days).", ACCENT_BLUE),
    ("\U0001F4CF", "Measure Roughness", "Measures how 'rough' price\nmovements are at MULTIPLE\ntime scales simultaneously.", PURPLE),
    ("\U0001F4CA", "Calculate Stress", "Same roughness everywhere = CALM\nDifferent roughness = STRESSED\nOutputs a score from 0 to 1.", ACCENT_AMBER),
    ("\U0001F50D", "Cross-Check", "Compares NIFTY with BANKNIFTY.\nIf they diverge, something\nunusual is happening.", ACCENT_GREEN),
    ("\U0001F916", "AI Explains", "AI reads all the data and\nexplains it in plain English\nwith actionable advice.", ACCENT_RED),
]

for i, (icon, title, desc, color) in enumerate(steps):
    x = Inches(0.6 + i * 2.5)
    y = Inches(2.35)
    # Card
    add_shape(slide, x, y, Inches(2.2), Inches(3.2), BG_CARD, color)
    # Step number
    num_badge = add_shape(slide, x + Inches(0.15), y + Inches(0.15), Inches(0.4), Inches(0.4), color)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.17), Inches(0.4), Inches(0.4),
                 str(i+1), font_size=16, color=BG_DARK, bold=True, alignment=PP_ALIGN.CENTER)
    # Icon
    add_text_box(slide, x, y + Inches(0.7), Inches(2.2), Inches(0.6),
                 icon, font_size=36, color=WHITE, alignment=PP_ALIGN.CENTER)
    # Title
    add_text_box(slide, x + Inches(0.15), y + Inches(1.35), Inches(1.9), Inches(0.4),
                 title, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Desc
    add_text_box(slide, x + Inches(0.15), y + Inches(1.8), Inches(1.9), Inches(1.2),
                 desc, font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # Arrow between cards (except last)
    if i < 4:
        add_text_box(slide, x + Inches(2.2), y + Inches(1.2), Inches(0.3), Inches(0.5),
                     "\u25B6", font_size=18, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom result bar
add_shape(slide, Inches(0.6), Inches(5.85), Inches(12.1), Inches(1.0), BG_CARD)
add_text_box(slide, Inches(0.9), Inches(5.95), Inches(3.5), Inches(0.35),
             "THE RESULT:  A Simple Traffic Light", font_size=16, color=WHITE, bold=True)

signals = [("GREEN", "Market is Calm", ACCENT_GREEN), ("AMBER", "Use Caution", ACCENT_AMBER), ("RED", "High Stress", ACCENT_RED)]
for i, (label, desc, color) in enumerate(signals):
    x = Inches(5.0 + i * 2.7)
    add_shape(slide, x, Inches(6.0), Inches(0.6), Inches(0.6), color)
    add_text_box(slide, x + Inches(0.75), Inches(5.95), Inches(1.8), Inches(0.35),
                 label, font_size=15, color=color, bold=True)
    add_text_box(slide, x + Inches(0.75), Inches(6.25), Inches(1.8), Inches(0.3),
                 desc, font_size=11, color=LIGHT_GRAY)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 4: THE SCIENCE - MFDFA
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, 3, "The Science: MFDFA", "Multifractal Detrended Fluctuation Analysis")

# Left - Simple explanation
add_shape(slide, Inches(0.6), Inches(1.7), Inches(5.8), Inches(5.2), BG_CARD, ACCENT_BLUE)
add_text_box(slide, Inches(0.9), Inches(1.85), Inches(5.2), Inches(0.4),
             "FOR EVERYONE: The Roughness Detector", font_size=16, color=ACCENT_BLUE, bold=True)
add_accent_line(slide, Inches(0.9), Inches(2.3), Inches(1.5), ACCENT_BLUE)

simple_lines = [
    ("Imagine looking at a mountain range:", 14, WHITE, True),
    ("", 6, WHITE, False),
    ("  \u2022  Smooth hills = Calm market (predictable)", 13, LIGHT_GRAY, False),
    ("  \u2022  Jagged peaks = Stressed market (chaotic)", 13, LIGHT_GRAY, False),
    ("", 8, WHITE, False),
    ("MFDFA looks at roughness at EVERY zoom level:", 14, WHITE, True),
    ("", 6, WHITE, False),
    ("  \U0001F50D  Zoom in:  5-50 minute moves", 13, ACCENT_GREEN, False),
    ("  \U0001F50E  Mid zoom:  1-4 hour moves", 13, ACCENT_AMBER, False),
    ("  \U0001F30D  Zoom out:  Multi-day swings", 13, PURPLE, False),
    ("", 8, WHITE, False),
    ("If roughness is SAME at all levels = CALM", 13, ACCENT_GREEN, True),
    ("If roughness DIFFERS across levels = STRESSED", 13, ACCENT_RED, True),
    ("", 8, WHITE, False),
    ("The wider the difference, the higher the stress score.", 13, WHITE, False),
]
add_multi_text(slide, Inches(0.9), Inches(2.5), Inches(5.2), Inches(4.2), simple_lines, line_spacing=1.1)

# Right - Technical
add_shape(slide, Inches(6.8), Inches(1.7), Inches(5.9), Inches(5.2), BG_CARD, PURPLE)
add_text_box(slide, Inches(7.1), Inches(1.85), Inches(5.3), Inches(0.4),
             "FOR EXPERTS: The Mathematics", font_size=16, color=PURPLE, bold=True)
add_accent_line(slide, Inches(7.1), Inches(2.3), Inches(1.5), PURPLE)

tech_lines = [
    ("Algorithm Steps:", 14, WHITE, True),
    ("", 4, WHITE, False),
    ("1. Profile:  Y(i) = \u03A3[P(k) - mean(P)]", 12, LIGHT_GRAY, False),
    ("2. Segment into scales s \u2208 [10, 250]", 12, LIGHT_GRAY, False),
    ("3. Local polynomial detrending (order m=1)", 12, LIGHT_GRAY, False),
    ("4. q-th order fluctuation: F_q(s)", 12, LIGHT_GRAY, False),
    ("5. Scaling:  F_q(s) ~ s^h(q)", 12, LIGHT_GRAY, False),
    ("", 6, WHITE, False),
    ("Key Outputs:", 14, WHITE, True),
    ("", 4, WHITE, False),
    ("\u2022  h(q) - Generalized Hurst exponents", 12, ACCENT_BLUE, False),
    ("\u2022  h(2) - Standard Hurst (>0.5 = trending)", 12, ACCENT_BLUE, False),
    ("\u2022  \u0394\u03B1 - Spectral width (wider = stressed)", 12, ACCENT_BLUE, False),
    ("\u2022  \u03C4(q) = q\u00B7h(q) - 1 (mass exponent)", 12, ACCENT_BLUE, False),
    ("\u2022  f(\u03B1) - Singularity spectrum", 12, ACCENT_BLUE, False),
    ("", 6, WHITE, False),
    ("q-orders: [-5, -3, -1, 0, 1, 2, 3, 5]", 12, MID_GRAY, False),
    ("Window: 1000 bars | Scales: 20 log-spaced", 12, MID_GRAY, False),
]
add_multi_text(slide, Inches(7.1), Inches(2.5), Inches(5.3), Inches(4.2), tech_lines, line_spacing=1.05)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 5: STRESS SCORE & REGIME
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, 4, "Stress Score & Regime Classification",
             "How raw math becomes a simple traffic light")

# Stress score formula card
add_shape(slide, Inches(0.6), Inches(1.7), Inches(7.5), Inches(2.4), BG_CARD, ACCENT_BLUE)
add_text_box(slide, Inches(0.9), Inches(1.85), Inches(6), Inches(0.4),
             "STRESS SCORE FORMULA", font_size=16, color=ACCENT_BLUE, bold=True)
add_accent_line(slide, Inches(0.9), Inches(2.3), Inches(1.2), ACCENT_BLUE)

formula_lines = [
    ("stress = 0.35 \u00D7 spectral_width  +  0.25 \u00D7 hurst_deviation", 15, WHITE, True),
    ("       + 0.20 \u00D7 asymmetry        +  0.20 \u00D7 hurst_volatility", 15, WHITE, True),
    ("", 8, WHITE, False),
    ("ensemble = 0.70 \u00D7 stress_score   +  0.30 \u00D7 coupling_score", 15, ACCENT_BLUE, True),
]
add_multi_text(slide, Inches(0.9), Inches(2.45), Inches(7), Inches(1.5), formula_lines, line_spacing=1.15)

# Component explanations
components = [
    ("35%", "Spectral Width", "How multifractal\nthe market is", ACCENT_BLUE),
    ("25%", "Hurst Deviation", "Distance of h(2)\nfrom random walk", PURPLE),
    ("20%", "Asymmetry", "Large vs small\nfluctuation scaling", ACCENT_AMBER),
    ("20%", "Hurst Volatility", "Inconsistency\nacross q-orders", ACCENT_RED),
]
for i, (pct, name, desc, color) in enumerate(components):
    x = Inches(0.6 + i * 2.0)
    y = Inches(4.35)
    add_shape(slide, x, y, Inches(1.8), Inches(1.6), BG_CARD, color)
    add_text_box(slide, x, y + Inches(0.1), Inches(1.8), Inches(0.4),
                 pct, font_size=22, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, y + Inches(0.55), Inches(1.8), Inches(0.35),
                 name, font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, y + Inches(0.9), Inches(1.8), Inches(0.6),
                 desc, font_size=10, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Regime classification on right
add_shape(slide, Inches(8.5), Inches(1.7), Inches(4.2), Inches(4.95), BG_CARD)
add_text_box(slide, Inches(8.8), Inches(1.85), Inches(3.6), Inches(0.4),
             "REGIME CLASSIFICATION", font_size=16, color=WHITE, bold=True)
add_accent_line(slide, Inches(8.8), Inches(2.3), Inches(1.2), ACCENT_GREEN)

regimes = [
    ("GREEN", "\u2264 0.35", "Calm, trending market\nTrade normally\nRisk: 1.5% per trade", ACCENT_GREEN),
    ("AMBER", "0.35 - 0.65", "Getting choppy\nReduce positions\nRisk: 1.0% per trade", ACCENT_AMBER),
    ("RED", "> 0.65", "High stress\nAvoid new trades\nRisk: 0.5% per trade", ACCENT_RED),
]
for i, (label, threshold, desc, color) in enumerate(regimes):
    y = Inches(2.55 + i * 1.3)
    add_shape(slide, Inches(8.8), y, Inches(3.6), Inches(1.15), fill_color=BG_DARK, border_color=color)
    # Color dot
    add_circle(slide, Inches(9.0), y + Inches(0.15), Inches(0.45), color)
    add_text_box(slide, Inches(9.6), y + Inches(0.08), Inches(1.2), Inches(0.3),
                 label, font_size=16, color=color, bold=True)
    add_text_box(slide, Inches(10.8), y + Inches(0.08), Inches(1.4), Inches(0.3),
                 threshold, font_size=13, color=MID_GRAY)
    add_text_box(slide, Inches(9.6), y + Inches(0.42), Inches(2.6), Inches(0.7),
                 desc, font_size=10, color=LIGHT_GRAY)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 6: FEATURES OVERVIEW
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, 5, "Features at a Glance", "12 powerful features in one platform")

features = [
    ("\U0001F6A6", "Stress Detection", "Traffic light regime\nclassification", ACCENT_RED),
    ("\U0001F4C8", "Trend Analysis", "EMA, RSI, MACD\nATR indicators", ACCENT_BLUE),
    ("\U0001F3AF", "Support & Resistance", "Fractal-validated\nprice levels", ACCENT_GREEN),
    ("\U0001F4A5", "Breakout Detection", "Quality-scored\nbreakout signals", ACCENT_AMBER),
    ("\U0001F6E1", "Risk Management", "Regime-aware stops\nand targets", PURPLE),
    ("\U0001F3C6", "Health Score", "Market grade\nfrom A+ to F", ACCENT_BLUE),
    ("\U0001F52E", "Regime Prediction", "Markov chain\ntransition forecast", ACCENT_GREEN),
    ("\U0001F680", "Topless Detection", "Price discovery\nidentification", ACCENT_RED),
    ("\U0001F4CA", "Multi-Stock", "70+ stocks\nanalysis on demand", ACCENT_AMBER),
    ("\U0001F916", "AI Insights", "Claude-powered\nplain English advice", PURPLE),
    ("\U0001F517", "Coupling Monitor", "Cross-asset\ncorrelation tracking", ACCENT_BLUE),
    ("\U0001F514", "Alert System", "Auto-triggered\nmarket alerts", ACCENT_GREEN),
]

for i, (icon, title, desc, color) in enumerate(features):
    row = i // 4
    col = i % 4
    x = Inches(0.6 + col * 3.15)
    y = Inches(1.75 + row * 1.85)
    card = add_shape(slide, x, y, Inches(2.85), Inches(1.6), BG_CARD, color)
    add_text_box(slide, x, y + Inches(0.1), Inches(2.85), Inches(0.5),
                 icon, font_size=28, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, y + Inches(0.6), Inches(2.85), Inches(0.35),
                 title, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, y + Inches(0.95), Inches(2.85), Inches(0.6),
                 desc, font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 7: ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, 6, "System Architecture", "How all the pieces fit together")

# Data Sources (bottom layer)
add_text_box(slide, Inches(0.6), Inches(1.55), Inches(3), Inches(0.35),
             "DATA SOURCES", font_size=13, color=ACCENT_GREEN, bold=True)
sources = [("Yahoo Finance", "Free, 15m delay"), ("Angel One API", "Real-time, 0 delay"), ("Anthropic Claude", "AI analysis")]
for i, (name, desc) in enumerate(sources):
    x = Inches(0.6 + i * 2.2)
    add_shape(slide, x, Inches(1.95), Inches(2.0), Inches(0.9), BG_CARD, ACCENT_GREEN)
    add_text_box(slide, x + Inches(0.1), Inches(2.0), Inches(1.8), Inches(0.35),
                 name, font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), Inches(2.35), Inches(1.8), Inches(0.35),
                 desc, font_size=10, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Arrow down
add_text_box(slide, Inches(3.2), Inches(2.9), Inches(0.5), Inches(0.4),
             "\u25BC", font_size=20, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

# Database layer
add_text_box(slide, Inches(0.6), Inches(3.25), Inches(3), Inches(0.35),
             "DATABASE", font_size=13, color=ACCENT_BLUE, bold=True)
add_shape(slide, Inches(0.6), Inches(3.6), Inches(6.4), Inches(0.7), BG_CARD, ACCENT_BLUE)
add_text_box(slide, Inches(0.9), Inches(3.65), Inches(5.8), Inches(0.25),
             "SQLite  |  candles  |  signals  |  ticks  |  alerts  |  portfolios  |  health",
             font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0.9), Inches(3.95), Inches(5.8), Inches(0.25),
             "Zero-config local storage  \u2022  Easily swappable to PostgreSQL/TimescaleDB",
             font_size=10, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Arrow down
add_text_box(slide, Inches(3.2), Inches(4.35), Inches(0.5), Inches(0.4),
             "\u25BC", font_size=20, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

# Compute Engines
add_text_box(slide, Inches(0.6), Inches(4.65), Inches(3), Inches(0.35),
             "COMPUTE ENGINES", font_size=13, color=PURPLE, bold=True)

engines = [
    ("MFDFA\nEngine", ACCENT_RED), ("Feature\nEngine", ACCENT_AMBER), ("Trend\nEngine", ACCENT_BLUE),
    ("S/R\nEngine", ACCENT_GREEN), ("Breakout\nEngine", PURPLE), ("Risk\nEngine", ACCENT_RED),
    ("Health\nScore", ACCENT_AMBER), ("Regime\nEngine", ACCENT_BLUE),
]
for i, (name, color) in enumerate(engines):
    x = Inches(0.6 + i * 1.58)
    add_shape(slide, x, Inches(5.0), Inches(1.38), Inches(0.85), BG_CARD, color)
    add_text_box(slide, x, Inches(5.05), Inches(1.38), Inches(0.75),
                 name, font_size=10, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Right side - API & UI
add_text_box(slide, Inches(7.4), Inches(1.55), Inches(3), Inches(0.35),
             "API LAYER", font_size=13, color=ACCENT_AMBER, bold=True)
add_shape(slide, Inches(7.4), Inches(1.95), Inches(5.3), Inches(2.1), BG_CARD, ACCENT_AMBER)

endpoints = [
    "/latest - Current regime & analysis",
    "/candles - OHLCV price history",
    "/signals - Signal history",
    "/spectrum - MFDFA spectrum",
    "/analysis - Full technical analysis",
    "/analyze/{symbol} - On-demand analysis",
    "/stocks - Available symbols (70+)",
    "/health - System status & alerts",
]
for i, ep in enumerate(endpoints):
    add_text_box(slide, Inches(7.6), Inches(2.05 + i * 0.24), Inches(4.9), Inches(0.25),
                 ep, font_size=10, color=LIGHT_GRAY, font_name="Consolas")

# UI
add_text_box(slide, Inches(7.4), Inches(4.3), Inches(3), Inches(0.35),
             "USER INTERFACE", font_size=13, color=ACCENT_BLUE, bold=True)
add_shape(slide, Inches(7.4), Inches(4.65), Inches(5.3), Inches(2.2), BG_CARD, ACCENT_BLUE)
ui_items = [
    "React Dashboard with Recharts",
    "Regime traffic light indicator",
    "Real-time stress score timeline",
    "Candlestick chart + EMA overlays",
    "RSI & MACD panels",
    "Support/Resistance levels table",
    "Risk management panel",
    "Multi-stock selector dropdown",
]
for i, item in enumerate(ui_items):
    add_text_box(slide, Inches(7.6), Inches(4.75 + i * 0.24), Inches(4.9), Inches(0.25),
                 "\u2022  " + item, font_size=11, color=LIGHT_GRAY)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 8: TECHNICAL ANALYSIS ENGINES
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, 7, "Technical Analysis Engines", "7 specialized engines working in concert")

engines_detail = [
    ("Trend Engine", "EMA (9/21/50), RSI (14), MACD (12/26/9), ATR (14)\nClassifies: STRONG_UP / UP / SIDEWAYS / DOWN / STRONG_DOWN", ACCENT_BLUE),
    ("S/R Engine", "Swing point clustering + Pivot points + Fractal validation\nRanks levels by touch count, groups within 0.3% tolerance", ACCENT_GREEN),
    ("Breakout Engine", "Volume (>1.5x avg) + Body (>1.2x avg) confirmation\nFractal quality score: STRONG_BUY / BUY / WATCH / AVOID", ACCENT_AMBER),
    ("Risk Engine", "ATR-based & S/R-based stoploss computation\nRegime-adjusted: GREEN 1.5% | AMBER 1.0% | RED 0.5% risk", ACCENT_RED),
    ("Topless Engine", "Detects price discovery (no overhead resistance)\nConditions: h(2) > 0.55, narrow spectrum, above-avg volume", PURPLE),
    ("Health Score", "Weighted composite: Regime (40%) + Hurst (20%) + Coupling (15%)\n+ RSI (10%) + Volatility (15%) = Grade from A+ to F", ACCENT_BLUE),
    ("Regime Engine", "Markov chain transition matrix from signal history\nLaplace smoothing + trend adjustment for predictions", ACCENT_GREEN),
]

for i, (name, desc, color) in enumerate(engines_detail):
    row = i // 2
    col = i % 2
    x = Inches(0.6 + col * 6.3)
    y = Inches(1.7 + row * 1.4)
    w = Inches(5.9) if i < 6 else Inches(12.2)
    add_shape(slide, x, y, w, Inches(1.2), BG_CARD, color)
    # Colored left bar
    add_shape(slide, x, y, Inches(0.08), Inches(1.2), color)
    add_text_box(slide, x + Inches(0.25), y + Inches(0.1), w - Inches(0.4), Inches(0.3),
                 name, font_size=15, color=color, bold=True)
    add_text_box(slide, x + Inches(0.25), y + Inches(0.45), w - Inches(0.4), Inches(0.7),
                 desc, font_size=12, color=LIGHT_GRAY)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 9: BREAKOUT QUALITY SCORING
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, 8, "Fractal-Aware Breakout Quality", "Not all breakouts are equal")

# Main formula
add_shape(slide, Inches(0.6), Inches(1.7), Inches(12.1), Inches(1.6), BG_CARD, ACCENT_AMBER)
add_text_box(slide, Inches(0.9), Inches(1.85), Inches(11.5), Inches(0.35),
             "BREAKOUT QUALITY FORMULA", font_size=16, color=ACCENT_AMBER, bold=True)

formula = ("quality = 0.25 \u00D7 hurst_persistence  +  0.15 \u00D7 spectrum_stability  +  "
           "0.20 \u00D7 mfdcca_alignment  +  0.15 \u00D7 sr_fractal_quality  +  0.25 \u00D7 classical_confirmation")
add_text_box(slide, Inches(0.9), Inches(2.35), Inches(11.5), Inches(0.4),
             formula, font_size=14, color=WHITE, bold=True, font_name="Consolas")
add_text_box(slide, Inches(0.9), Inches(2.75), Inches(11.5), Inches(0.3),
             "Combines fractal mathematics with classical volume/body confirmation for robust signal quality",
             font_size=12, color=MID_GRAY)

# 5 component cards
comps = [
    ("25%", "Hurst\nPersistence", "h(2) supporting\ntrend continuation", ACCENT_BLUE),
    ("15%", "Spectrum\nStability", "Narrow spectrum =\nclean breakout", PURPLE),
    ("20%", "MFDCCA\nAlignment", "Cross-asset\ncoupling supports", ACCENT_GREEN),
    ("15%", "S/R Fractal\nQuality", "Level validated by\nfractal scales", ACCENT_AMBER),
    ("25%", "Classical\nConfirmation", "Volume > 1.5x avg\nBody > 1.2x avg", ACCENT_RED),
]
for i, (pct, name, desc, color) in enumerate(comps):
    x = Inches(0.6 + i * 2.5)
    y = Inches(3.65)
    add_shape(slide, x, y, Inches(2.2), Inches(2.0), BG_CARD, color)
    add_text_box(slide, x, y + Inches(0.1), Inches(2.2), Inches(0.45),
                 pct, font_size=26, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, y + Inches(0.6), Inches(2.2), Inches(0.5),
                 name, font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, y + Inches(1.2), Inches(2.2), Inches(0.7),
                 desc, font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Signal labels at bottom
add_shape(slide, Inches(0.6), Inches(5.95), Inches(12.1), Inches(0.9), BG_CARD)
add_text_box(slide, Inches(0.9), Inches(6.05), Inches(2), Inches(0.3),
             "SIGNAL LABELS:", font_size=14, color=WHITE, bold=True)
labels = [
    ("STRONG_BUY", "Quality > 0.8", ACCENT_GREEN),
    ("BUY", "Quality 0.6 - 0.8", ACCENT_BLUE),
    ("WATCH", "Quality 0.4 - 0.6", ACCENT_AMBER),
    ("AVOID", "Quality < 0.4", ACCENT_RED),
]
for i, (label, rng, color) in enumerate(labels):
    x = Inches(3.2 + i * 2.5)
    add_shape(slide, x, Inches(6.05), Inches(0.2), Inches(0.65), color)
    add_text_box(slide, x + Inches(0.35), Inches(6.05), Inches(1.8), Inches(0.3),
                 label, font_size=13, color=color, bold=True)
    add_text_box(slide, x + Inches(0.35), Inches(6.35), Inches(1.8), Inches(0.3),
                 rng, font_size=11, color=LIGHT_GRAY)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 10: AI-POWERED INSIGHTS
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, 9, "AI-Powered Market Insights", "Claude AI translates math into actionable advice")

# Left - How it works
add_shape(slide, Inches(0.6), Inches(1.7), Inches(5.8), Inches(5.2), BG_CARD, PURPLE)
add_text_box(slide, Inches(0.9), Inches(1.85), Inches(5.2), Inches(0.4),
             "HOW THE AI WORKS", font_size=16, color=PURPLE, bold=True)
add_accent_line(slide, Inches(0.9), Inches(2.3), Inches(1.2), PURPLE)

steps_ai = [
    ("\u2460  Context Building", "All computed data (regime, stress, trend, S/R,\nbreakout, risk, health) compiled into a\nstructured prompt for the AI.", WHITE),
    ("\u2461  Analysis Generation", "Claude produces a detailed market report\ncovering current state, key levels,\nrisk assessment, and recommendations.", WHITE),
    ("\u2462  Interactive Chat", "Traders can ask follow-up questions:\n\u2022 \"Should I go long here?\"\n\u2022 \"What's the risk of holding overnight?\"\n\u2022 \"Explain the current Hurst exponent\"", WHITE),
]
for i, (title, desc, _) in enumerate(steps_ai):
    y = Inches(2.55 + i * 1.55)
    add_text_box(slide, Inches(0.9), y, Inches(5.2), Inches(0.3),
                 title, font_size=14, color=PURPLE, bold=True)
    add_text_box(slide, Inches(0.9), y + Inches(0.35), Inches(5.2), Inches(1.1),
                 desc, font_size=12, color=LIGHT_GRAY)

# Right - Example output
add_shape(slide, Inches(6.8), Inches(1.7), Inches(5.9), Inches(5.2), RGBColor(0x0D, 0x1A, 0x0D), ACCENT_GREEN)
add_text_box(slide, Inches(7.1), Inches(1.85), Inches(5.3), Inches(0.4),
             "EXAMPLE AI OUTPUT", font_size=16, color=ACCENT_GREEN, bold=True)
add_accent_line(slide, Inches(7.1), Inches(2.3), Inches(1.2), ACCENT_GREEN)

output_lines = [
    ("MARKET ANALYSIS - NIFTY 50", 13, ACCENT_GREEN, True),
    ("", 6, WHITE, False),
    ("Current State: GREEN regime (Score: 0.28)", 12, WHITE, True),
    ("Market is calm and trending. Hurst at 0.62", 11, LIGHT_GRAY, False),
    ("indicates strong persistence.", 11, LIGHT_GRAY, False),
    ("", 6, WHITE, False),
    ("Key Levels:", 12, WHITE, True),
    ("  Resistance: 24,520 (3 touches, validated)", 11, ACCENT_RED, False),
    ("  Support: 24,280 (5 touches, strong)", 11, ACCENT_GREEN, False),
    ("  Pivot: 24,400", 11, ACCENT_AMBER, False),
    ("", 6, WHITE, False),
    ("Breakout: Approaching 24,520", 12, WHITE, True),
    ("Quality: BUY (0.72). Volume confirming.", 11, ACCENT_BLUE, False),
    ("", 6, WHITE, False),
    ("Risk Management:", 12, WHITE, True),
    ("  Entry: Above 24,530", 11, LIGHT_GRAY, False),
    ("  Stoploss: 24,380 (ATR-based, 150 pts)", 11, ACCENT_RED, False),
    ("  Target 1: 24,680 | Target 2: 24,830", 11, ACCENT_GREEN, False),
    ("", 6, WHITE, False),
    ("Health Grade: A (84/100)", 12, ACCENT_GREEN, True),
]
add_multi_text(slide, Inches(7.1), Inches(2.5), Inches(5.3), Inches(4.2), output_lines, line_spacing=1.0)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 11: DATA SOURCES
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, 10, "Data Sources & Integration", "Flexible data pipeline with free and premium options")

# Yahoo Finance card
add_shape(slide, Inches(0.6), Inches(1.7), Inches(3.8), Inches(4.5), BG_CARD, ACCENT_BLUE)
add_text_box(slide, Inches(0.9), Inches(1.85), Inches(3.2), Inches(0.4),
             "YAHOO FINANCE", font_size=18, color=ACCENT_BLUE, bold=True)
add_text_box(slide, Inches(0.9), Inches(2.3), Inches(3.2), Inches(0.3),
             "DEFAULT  \u2022  FREE", font_size=12, color=ACCENT_GREEN, bold=True)
yahoo_info = [
    ("Cost:", "Free", WHITE, ACCENT_GREEN),
    ("Delay:", "~15 minutes", WHITE, ACCENT_AMBER),
    ("History:", "60 days of 5-min candles", WHITE, LIGHT_GRAY),
    ("Auth:", "None required", WHITE, LIGHT_GRAY),
    ("Coverage:", "NIFTY 50, BANKNIFTY,\n50+ stocks, globals", WHITE, LIGHT_GRAY),
    ("Best For:", "Swing trading, learning,\nend-of-day analysis", WHITE, LIGHT_GRAY),
]
for i, (label, value, lc, vc) in enumerate(yahoo_info):
    y = Inches(2.75 + i * 0.5)
    add_text_box(slide, Inches(0.9), y, Inches(1.2), Inches(0.3), label, font_size=11, color=MID_GRAY, bold=True)
    add_text_box(slide, Inches(2.1), y, Inches(2.1), Inches(0.45), value, font_size=11, color=vc)

# Angel One card
add_shape(slide, Inches(4.75), Inches(1.7), Inches(3.8), Inches(4.5), BG_CARD, ACCENT_AMBER)
add_text_box(slide, Inches(5.05), Inches(1.85), Inches(3.2), Inches(0.4),
             "ANGEL ONE SmartAPI", font_size=18, color=ACCENT_AMBER, bold=True)
add_text_box(slide, Inches(5.05), Inches(2.3), Inches(3.2), Inches(0.3),
             "OPTIONAL  \u2022  REAL-TIME", font_size=12, color=ACCENT_AMBER, bold=True)
angel_info = [
    ("Cost:", "Free with account", WHITE, ACCENT_GREEN),
    ("Delay:", "0 (real-time)", WHITE, ACCENT_GREEN),
    ("Auth:", "Client ID + TOTP", WHITE, LIGHT_GRAY),
    ("Rate:", "5 requests/second", WHITE, LIGHT_GRAY),
    ("Coverage:", "All NSE/BSE listed\nsecurities", WHITE, LIGHT_GRAY),
    ("Best For:", "Intraday trading,\nreal-time monitoring", WHITE, LIGHT_GRAY),
]
for i, (label, value, lc, vc) in enumerate(angel_info):
    y = Inches(2.75 + i * 0.5)
    add_text_box(slide, Inches(5.05), y, Inches(1.2), Inches(0.3), label, font_size=11, color=MID_GRAY, bold=True)
    add_text_box(slide, Inches(6.25), y, Inches(2.1), Inches(0.45), value, font_size=11, color=vc)

# Claude AI card
add_shape(slide, Inches(8.9), Inches(1.7), Inches(3.8), Inches(4.5), BG_CARD, PURPLE)
add_text_box(slide, Inches(9.2), Inches(1.85), Inches(3.2), Inches(0.4),
             "ANTHROPIC CLAUDE", font_size=18, color=PURPLE, bold=True)
add_text_box(slide, Inches(9.2), Inches(2.3), Inches(3.2), Inches(0.3),
             "AI INSIGHTS  \u2022  OPTIONAL", font_size=12, color=PURPLE, bold=True)
claude_info = [
    ("Model:", "Claude Sonnet", WHITE, PURPLE),
    ("Purpose:", "Market analysis +\ninteractive Q&A", WHITE, LIGHT_GRAY),
    ("Auth:", "API key required", WHITE, LIGHT_GRAY),
    ("Cost:", "Pay per token\n(very affordable)", WHITE, ACCENT_GREEN),
    ("Features:", "Analysis reports,\nchat interface", WHITE, LIGHT_GRAY),
    ("Best For:", "Understanding complex\nfractal signals easily", WHITE, LIGHT_GRAY),
]
for i, (label, value, lc, vc) in enumerate(claude_info):
    y = Inches(2.75 + i * 0.5)
    add_text_box(slide, Inches(9.2), y, Inches(1.2), Inches(0.3), label, font_size=11, color=MID_GRAY, bold=True)
    add_text_box(slide, Inches(10.4), y, Inches(2.1), Inches(0.45), value, font_size=11, color=vc)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 12: DASHBOARD WIREFRAME
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, 11, "The Dashboard", "Interactive web interface with real-time charts")

# Mock dashboard layout
# Top bar
add_shape(slide, Inches(0.6), Inches(1.65), Inches(12.1), Inches(0.6), RGBColor(0x10, 0x1B, 0x30))
add_text_box(slide, Inches(0.9), Inches(1.7), Inches(4), Inches(0.5),
             "NIFTY Stress Monitor", font_size=16, color=WHITE, bold=True)
add_text_box(slide, Inches(9.5), Inches(1.7), Inches(3), Inches(0.5),
             "[Stock Selector \u25BC]  [Refresh]", font_size=12, color=MID_GRAY)

# Regime + Health row
add_shape(slide, Inches(0.6), Inches(2.4), Inches(4.0), Inches(1.2), BG_CARD, ACCENT_GREEN)
add_circle(slide, Inches(1.0), Inches(2.6), Inches(0.7), ACCENT_GREEN)
add_text_box(slide, Inches(1.9), Inches(2.55), Inches(2.5), Inches(0.35),
             "GREEN", font_size=22, color=ACCENT_GREEN, bold=True)
add_text_box(slide, Inches(1.9), Inches(2.9), Inches(2.5), Inches(0.5),
             "Ensemble: 0.28 | Stress: 0.25", font_size=11, color=LIGHT_GRAY)

add_shape(slide, Inches(4.85), Inches(2.4), Inches(3.5), Inches(1.2), BG_CARD, ACCENT_BLUE)
add_text_box(slide, Inches(5.15), Inches(2.55), Inches(1.5), Inches(0.4),
             "84", font_size=30, color=ACCENT_BLUE, bold=True)
add_text_box(slide, Inches(6.4), Inches(2.55), Inches(1.7), Inches(0.35),
             "/ 100  Grade: A", font_size=14, color=WHITE)
add_text_box(slide, Inches(5.15), Inches(3.0), Inches(3.0), Inches(0.3),
             "Health Score - Very Good", font_size=11, color=LIGHT_GRAY)

add_shape(slide, Inches(8.6), Inches(2.4), Inches(4.1), Inches(1.2), BG_CARD, ACCENT_AMBER)
add_text_box(slide, Inches(8.9), Inches(2.55), Inches(3.5), Inches(0.35),
             "BREAKOUT: BULLISH", font_size=14, color=ACCENT_AMBER, bold=True)
add_text_box(slide, Inches(8.9), Inches(2.9), Inches(3.5), Inches(0.5),
             "Quality: BUY (0.72) | Vol: OK", font_size=11, color=LIGHT_GRAY)

# Stress timeline mock
add_shape(slide, Inches(0.6), Inches(3.8), Inches(7.8), Inches(1.5), BG_CARD)
add_text_box(slide, Inches(0.9), Inches(3.85), Inches(3), Inches(0.3),
             "STRESS SCORE TIMELINE", font_size=11, color=MID_GRAY, bold=True)
# Mock chart line
for i in range(30):
    x = Inches(1.0 + i * 0.24)
    h_val = 0.15 + (i % 7) * 0.03 + (0.02 if i > 20 else 0)
    bar_h = Inches(h_val * 3)
    color = ACCENT_GREEN if h_val < 0.28 else (ACCENT_AMBER if h_val < 0.38 else ACCENT_RED)
    add_shape(slide, x, Inches(5.1) - bar_h, Inches(0.15), bar_h, color)

# Price chart mock
add_shape(slide, Inches(8.6), Inches(3.8), Inches(4.1), Inches(1.5), BG_CARD)
add_text_box(slide, Inches(8.9), Inches(3.85), Inches(3), Inches(0.3),
             "PRICE CHART + EMAs", font_size=11, color=MID_GRAY, bold=True)

# S/R and Risk panels
add_shape(slide, Inches(0.6), Inches(5.5), Inches(4.0), Inches(1.5), BG_CARD)
add_text_box(slide, Inches(0.9), Inches(5.55), Inches(3.4), Inches(0.3),
             "SUPPORT / RESISTANCE", font_size=11, color=MID_GRAY, bold=True)
sr_items = ["R2: 24,650 (2 touches)", "R1: 24,520 (3 touches)", "PP: 24,400",
            "S1: 24,280 (5 touches)", "S2: 24,100 (1 touch)"]
for i, item in enumerate(sr_items):
    color = ACCENT_RED if "R" in item[:2] else (ACCENT_AMBER if "PP" in item else ACCENT_GREEN)
    add_text_box(slide, Inches(0.9), Inches(5.85 + i * 0.22), Inches(3.4), Inches(0.22),
                 item, font_size=10, color=color)

add_shape(slide, Inches(4.85), Inches(5.5), Inches(3.75), Inches(1.5), BG_CARD)
add_text_box(slide, Inches(5.15), Inches(5.55), Inches(3.1), Inches(0.3),
             "RISK MANAGEMENT", font_size=11, color=MID_GRAY, bold=True)
risk_items = [
    ("Long SL: 24,380", ACCENT_RED), ("Target 1: 24,680", ACCENT_GREEN),
    ("Target 2: 24,830", ACCENT_GREEN), ("R:R = 1:2.0", ACCENT_BLUE),
    ("Risk/Trade: 1.5%", WHITE),
]
for i, (item, color) in enumerate(risk_items):
    add_text_box(slide, Inches(5.15), Inches(5.85 + i * 0.22), Inches(3.1), Inches(0.22),
                 item, font_size=10, color=color)

# RSI / MACD mock
add_shape(slide, Inches(8.6), Inches(5.5), Inches(4.1), Inches(1.5), BG_CARD)
add_text_box(slide, Inches(8.9), Inches(5.55), Inches(3.5), Inches(0.3),
             "RSI & MACD PANELS", font_size=11, color=MID_GRAY, bold=True)
add_text_box(slide, Inches(8.9), Inches(5.9), Inches(3.5), Inches(0.8),
             "RSI: 58 (Neutral zone)\nMACD: Bullish crossover\nATR: 145 points\nTrend: UP (EMA9 > EMA21 > EMA50)",
             font_size=10, color=LIGHT_GRAY)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 13: ALERT SYSTEM
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, 12, "Alert System", "Automated notifications for critical market events")

alerts_data = [
    ("\U0001F6A8", "Regime Change", "HIGH", "Triggers when market regime transitions\n(e.g., GREEN \u2192 RED)",
     "NIFTY regime changed from GREEN to RED.\nConsider reducing positions and tightening stops.", ACCENT_RED),
    ("\U0001F4C8", "Stress Spike", "MEDIUM", "Triggers when ensemble score jumps\n>0.15 in a single signal",
     "Stress spike detected: 0.32 \u2192 0.52.\nMarket volatility increasing rapidly.", ACCENT_AMBER),
    ("\U0001F4A5", "Breakout", "MEDIUM", "Triggers when price breaks S/R level\nwith quality \u2265 BUY",
     "Bullish breakout at 24,520.\nQuality: BUY (0.72). Volume confirmed.", ACCENT_BLUE),
    ("\U0001F517", "Coupling Divergence", "LOW", "Triggers when coupling drops >0.2\nover 5 consecutive signals",
     "NIFTY-BANKNIFTY coupling diverging.\nSector rotation may be occurring.", PURPLE),
]

for i, (icon, name, severity, condition, example, color) in enumerate(alerts_data):
    y = Inches(1.7 + i * 1.38)
    # Main card
    add_shape(slide, Inches(0.6), y, Inches(12.1), Inches(1.2), BG_CARD, color)
    # Left color bar
    add_shape(slide, Inches(0.6), y, Inches(0.08), Inches(1.2), color)
    # Icon + name
    add_text_box(slide, Inches(0.9), y + Inches(0.12), Inches(0.5), Inches(0.4),
                 icon, font_size=24, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.5), y + Inches(0.12), Inches(2.2), Inches(0.35),
                 name, font_size=16, color=WHITE, bold=True)
    # Severity badge
    sev_color = ACCENT_RED if severity == "HIGH" else (ACCENT_AMBER if severity == "MEDIUM" else ACCENT_BLUE)
    badge = add_shape(slide, Inches(1.5), y + Inches(0.55), Inches(0.9), Inches(0.3), sev_color)
    add_text_box(slide, Inches(1.5), y + Inches(0.55), Inches(0.9), Inches(0.3),
                 severity, font_size=9, color=BG_DARK, bold=True, alignment=PP_ALIGN.CENTER)
    # Condition
    add_text_box(slide, Inches(3.8), y + Inches(0.12), Inches(3.5), Inches(0.9),
                 condition, font_size=11, color=LIGHT_GRAY)
    # Example
    add_text_box(slide, Inches(7.5), y + Inches(0.05), Inches(1.0), Inches(0.3),
                 "Example:", font_size=10, color=MID_GRAY, bold=True)
    add_text_box(slide, Inches(7.5), y + Inches(0.3), Inches(5.0), Inches(0.8),
                 example, font_size=11, color=color)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 14: TECH STACK
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, 13, "Tech Stack", "Lightweight, powerful, minimal dependencies")

# Backend
add_text_box(slide, Inches(0.6), Inches(1.6), Inches(2), Inches(0.35),
             "BACKEND", font_size=14, color=ACCENT_BLUE, bold=True)
backend_items = [
    ("Python 3.x", "Primary language", ACCENT_BLUE),
    ("FastAPI", "High-performance REST API", ACCENT_BLUE),
    ("NumPy", "Pure numerical computing (no pandas!)", ACCENT_BLUE),
    ("SQLite", "Zero-config local database", ACCENT_BLUE),
    ("Uvicorn", "ASGI production server", ACCENT_BLUE),
]
for i, (name, desc, color) in enumerate(backend_items):
    y = Inches(2.0 + i * 0.55)
    add_shape(slide, Inches(0.6), y, Inches(5.8), Inches(0.45), BG_CARD)
    add_text_box(slide, Inches(0.9), y + Inches(0.07), Inches(1.8), Inches(0.3),
                 name, font_size=13, color=color, bold=True)
    add_text_box(slide, Inches(2.8), y + Inches(0.07), Inches(3.4), Inches(0.3),
                 desc, font_size=12, color=LIGHT_GRAY)

# Frontend
add_text_box(slide, Inches(0.6), Inches(4.85), Inches(2), Inches(0.35),
             "FRONTEND", font_size=14, color=ACCENT_GREEN, bold=True)
frontend_items = [
    ("React", "Component-based UI framework", ACCENT_GREEN),
    ("Recharts", "Beautiful financial charts", ACCENT_GREEN),
    ("HTML5", "Pre-compiled single-file delivery", ACCENT_GREEN),
]
for i, (name, desc, color) in enumerate(frontend_items):
    y = Inches(5.25 + i * 0.55)
    add_shape(slide, Inches(0.6), y, Inches(5.8), Inches(0.45), BG_CARD)
    add_text_box(slide, Inches(0.9), y + Inches(0.07), Inches(1.8), Inches(0.3),
                 name, font_size=13, color=color, bold=True)
    add_text_box(slide, Inches(2.8), y + Inches(0.07), Inches(3.4), Inches(0.3),
                 desc, font_size=12, color=LIGHT_GRAY)

# Right side - Design decisions
add_shape(slide, Inches(6.8), Inches(1.6), Inches(5.9), Inches(5.3), BG_CARD, ACCENT_AMBER)
add_text_box(slide, Inches(7.1), Inches(1.75), Inches(5.3), Inches(0.4),
             "KEY DESIGN DECISIONS", font_size=16, color=ACCENT_AMBER, bold=True)
add_accent_line(slide, Inches(7.1), Inches(2.2), Inches(1.5), ACCENT_AMBER)

decisions = [
    ("No pandas", "Pure NumPy for all computation.\nFaster, lighter, fewer dependencies."),
    ("SQLite over Postgres", "Single-user app. No need for\nmulti-connection overhead."),
    ("Compiled HTML", "JSX pre-compiled to single file.\nNo separate frontend server."),
    ("Dual data sources", "Yahoo (free) for learning.\nAngel One (real-time) for trading."),
    ("4 core packages", "Minimal dependencies.\nRuns on any machine easily."),
    ("~6,100 lines total", "Compact but comprehensive.\n24 Python files + 1 React component."),
]
for i, (title, desc) in enumerate(decisions):
    y = Inches(2.45 + i * 0.72)
    add_circle(slide, Inches(7.2), y + Inches(0.05), Inches(0.3), ACCENT_AMBER)
    add_text_box(slide, Inches(7.2), y + Inches(0.02), Inches(0.3), Inches(0.3),
                 str(i+1), font_size=11, color=BG_DARK, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(7.65), y, Inches(2.0), Inches(0.3),
                 title, font_size=13, color=WHITE, bold=True)
    add_text_box(slide, Inches(9.7), y, Inches(2.7), Inches(0.6),
                 desc, font_size=11, color=LIGHT_GRAY)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 15: HOW TO RUN
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, 14, "Getting Started", "Up and running in 3 commands")

# Quick start
add_shape(slide, Inches(0.6), Inches(1.7), Inches(7.5), Inches(3.0), RGBColor(0x0D, 0x1A, 0x0D), ACCENT_GREEN)
add_text_box(slide, Inches(0.9), Inches(1.85), Inches(6), Inches(0.4),
             "QUICK START", font_size=18, color=ACCENT_GREEN, bold=True)
add_accent_line(slide, Inches(0.9), Inches(2.3), Inches(1.2), ACCENT_GREEN)

commands = [
    ("# 1. Install dependencies", MID_GRAY),
    ("pip install -r requirements.txt", ACCENT_GREEN),
    ("", WHITE),
    ("# 2. Run in DEMO mode (no internet needed)", MID_GRAY),
    ("python run.py demo", ACCENT_GREEN),
    ("", WHITE),
    ("# 3. Run in LIVE mode (real market data)", MID_GRAY),
    ("python run.py live", ACCENT_GREEN),
    ("", WHITE),
    ("# Open dashboard at http://localhost:8000", MID_GRAY),
]
for i, (cmd, color) in enumerate(commands):
    add_text_box(slide, Inches(0.9), Inches(2.5 + i * 0.28), Inches(6.8), Inches(0.28),
                 cmd, font_size=13, color=color, font_name="Consolas")

# Run modes
add_shape(slide, Inches(0.6), Inches(4.95), Inches(7.5), Inches(2.0), BG_CARD)
add_text_box(slide, Inches(0.9), Inches(5.1), Inches(6), Inches(0.35),
             "RUN MODES", font_size=14, color=WHITE, bold=True)

modes = [
    ("demo", "No", "Synthetic data", "Testing, demos, learning", ACCENT_GREEN),
    ("live", "Yes", "Yahoo Finance", "Actual trading analysis", ACCENT_BLUE),
    ("backfill", "Yes", "Yahoo Finance", "Pre-loading historical data", ACCENT_AMBER),
]
# Header
for j, header in enumerate(["Mode", "Internet?", "Data Source", "Best For"]):
    x = Inches(0.9 + j * 1.8)
    add_text_box(slide, x, Inches(5.5), Inches(1.6), Inches(0.3),
                 header, font_size=11, color=MID_GRAY, bold=True)
# Rows
for i, (mode, internet, source, best, color) in enumerate(modes):
    y = Inches(5.8 + i * 0.35)
    vals = [mode, internet, source, best]
    for j, val in enumerate(vals):
        x = Inches(0.9 + j * 1.8)
        c = color if j == 0 else LIGHT_GRAY
        add_text_box(slide, x, y, Inches(1.6), Inches(0.3),
                     val, font_size=11, color=c, bold=(j==0))

# Right - Prerequisites and env vars
add_shape(slide, Inches(8.5), Inches(1.7), Inches(4.2), Inches(2.3), BG_CARD, ACCENT_BLUE)
add_text_box(slide, Inches(8.8), Inches(1.85), Inches(3.6), Inches(0.35),
             "PREREQUISITES", font_size=14, color=ACCENT_BLUE, bold=True)
prereqs = ["Python 3.8+", "pip (package manager)", "Internet (for live mode only)"]
for i, p in enumerate(prereqs):
    add_text_box(slide, Inches(8.8), Inches(2.3 + i * 0.3), Inches(3.6), Inches(0.3),
                 "\u2713  " + p, font_size=12, color=LIGHT_GRAY)

add_shape(slide, Inches(8.5), Inches(4.2), Inches(4.2), Inches(2.75), BG_CARD, PURPLE)
add_text_box(slide, Inches(8.8), Inches(4.35), Inches(3.6), Inches(0.35),
             "OPTIONAL API KEYS (.env)", font_size=14, color=PURPLE, bold=True)
env_vars = [
    "ANTHROPIC_API_KEY",
    "  \u2192 AI-powered insights & chat",
    "",
    "ANGEL_API_KEY",
    "ANGEL_CLIENT_ID",
    "ANGEL_PIN",
    "ANGEL_TOTP_SECRET",
    "  \u2192 Real-time Angel One data",
]
for i, var in enumerate(env_vars):
    color = PURPLE if var.startswith("A") else (MID_GRAY if "\u2192" in var else LIGHT_GRAY)
    add_text_box(slide, Inches(8.8), Inches(4.8 + i * 0.28), Inches(3.6), Inches(0.28),
                 var, font_size=11, color=color, font_name="Consolas" if var.startswith("A") else "Segoe UI")


# ══════════════════════════════════════════════════════════════════════
# SLIDE 16: WHAT MAKES IT UNIQUE
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, 15, "What Makes It Unique", "7 differentiators from traditional trading tools")

differentiators = [
    ("\U0001F9EE", "Fractal Math in Retail", "Most retail platforms use basic TA.\nThis brings institutional-grade\nmultifractal analysis to everyone.", ACCENT_BLUE),
    ("\U0001F4D0", "Pure NumPy MFDFA", "No black boxes. Every calculation\nis transparent, inspectable, and\neducational. No TA-Lib dependency.", PURPLE),
    ("\U0001F504", "Regime-Aware Everything", "Stops, targets, position sizes,\nand breakout quality ALL adapt\nto the current market regime.", ACCENT_GREEN),
    ("\U0001F50D", "Multi-Scale Analysis", "20 simultaneous scales reveal\npatterns invisible to\nsingle-timeframe analysis.", ACCENT_AMBER),
    ("\U0001F517", "Cross-Asset Intelligence", "MFDCCA coupling detects sector\ndivergence before major moves.\nEarly warning system.", ACCENT_RED),
    ("\U0001F916", "AI + Math Fusion", "Raw fractal math explained in\nplain English by Claude AI.\nBest of both worlds.", PURPLE),
    ("\U0001F4B0", "Free & Private", "Free data, no subscriptions.\nRuns locally. Your data\nnever leaves your machine.", ACCENT_GREEN),
]

for i, (icon, title, desc, color) in enumerate(differentiators):
    row = i // 4
    col = i % 4
    if i == 6:  # Center the last one
        col = 1.5
    x = Inches(0.6 + col * 3.15)
    y = Inches(1.65 + row * 2.7)
    # Card
    add_shape(slide, x, y, Inches(2.85), Inches(2.45), BG_CARD, color)
    # Icon
    add_text_box(slide, x, y + Inches(0.15), Inches(2.85), Inches(0.5),
                 icon, font_size=32, alignment=PP_ALIGN.CENTER)
    # Title
    add_text_box(slide, x + Inches(0.15), y + Inches(0.7), Inches(2.55), Inches(0.4),
                 title, font_size=15, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    # Description
    add_text_box(slide, x + Inches(0.15), y + Inches(1.15), Inches(2.55), Inches(1.2),
                 desc, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 17: FUTURE SCOPE
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, 16, "Future Scope", "Where we can take this next")

future_items = [
    ("Options Chain Analysis", "IV surface analysis integrated\nwith fractal regime data", "Medium", ACCENT_BLUE),
    ("Multi-Timeframe MFDFA", "Simultaneous 1-min, 15-min,\nhourly fractal analysis", "Medium", PURPLE),
    ("Portfolio Optimization", "Fractal-aware allocation using\ncoupling matrix across assets", "Hard", ACCENT_RED),
    ("Backtesting Engine", "Test trading strategies against\nhistorical regime data", "Medium", ACCENT_GREEN),
    ("Push Notifications", "Alerts via Telegram, WhatsApp,\nor SMS in real-time", "Easy", ACCENT_GREEN),
    ("WebSocket Streaming", "Real-time dashboard updates\nwithout polling", "Medium", ACCENT_AMBER),
    ("Sector Rotation Map", "MFDCCA correlation matrix\nacross all market sectors", "Hard", ACCENT_RED),
    ("ML Regime Prediction", "Train classifier on MFDFA\nfeatures for forecasting", "Hard", PURPLE),
]

for i, (title, desc, difficulty, color) in enumerate(future_items):
    row = i // 4
    col = i % 4
    x = Inches(0.6 + col * 3.15)
    y = Inches(1.65 + row * 2.7)
    add_shape(slide, x, y, Inches(2.85), Inches(2.4), BG_CARD, color)
    # Difficulty badge
    diff_color = ACCENT_GREEN if difficulty == "Easy" else (ACCENT_AMBER if difficulty == "Medium" else ACCENT_RED)
    badge = add_shape(slide, x + Inches(0.15), y + Inches(0.15), Inches(0.9), Inches(0.3), diff_color)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.15), Inches(0.9), Inches(0.3),
                 difficulty, font_size=9, color=BG_DARK, bold=True, alignment=PP_ALIGN.CENTER)
    # Title
    add_text_box(slide, x + Inches(0.15), y + Inches(0.6), Inches(2.55), Inches(0.4),
                 title, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Description
    add_text_box(slide, x + Inches(0.15), y + Inches(1.1), Inches(2.55), Inches(1.1),
                 desc, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 18: CLOSING / THANK YOU
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

# Decorative circles
add_circle(slide, Inches(-1.5), Inches(-1), Inches(5), RGBColor(0x1A, 0x25, 0x3C))
add_circle(slide, Inches(9.5), Inches(4), Inches(6), RGBColor(0x15, 0x20, 0x35))

# Accent line
add_accent_line(slide, Inches(4.5), Inches(1.8), Inches(4.3), ACCENT_BLUE)

# Title
add_text_box(slide, Inches(1), Inches(2.1), Inches(11.3), Inches(1),
             "NIFTY Stress Monitor", font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Tagline
add_text_box(slide, Inches(1), Inches(3.2), Inches(11.3), Inches(0.6),
             "Where Fractal Mathematics Meets Practical Trading",
             font_size=22, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

# Quote
add_text_box(slide, Inches(2), Inches(4.2), Inches(9.3), Inches(0.8),
             "\"Markets are not random walks. They are fractals.\nAnd fractals can be measured, understood, and acted upon.\"",
             font_size=16, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Stats row
stats = [
    ("~6,100", "Lines of Code"),
    ("24", "Python Modules"),
    ("70+", "Supported Stocks"),
    ("12", "Analysis Engines"),
    ("0", "Black Boxes"),
]
for i, (num, label) in enumerate(stats):
    x = Inches(1.0 + i * 2.3)
    y = Inches(5.3)
    add_text_box(slide, x, y, Inches(2.0), Inches(0.5),
                 num, font_size=30, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, y + Inches(0.5), Inches(2.0), Inches(0.3),
                 label, font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom
add_text_box(slide, Inches(1), Inches(6.5), Inches(11.3), Inches(0.4),
             "Built with Python  \u2022  FastAPI  \u2022  NumPy  \u2022  React  \u2022  Recharts  \u2022  Claude AI",
             font_size=13, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(6.9), Inches(11.3), Inches(0.4),
             "Thank You",
             font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)


# ── Save ──────────────────────────────────────────────────────────────
output_path = os.path.join(os.path.dirname(__file__), "NIFTY_Stress_Monitor_Presentation.pptx")
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
